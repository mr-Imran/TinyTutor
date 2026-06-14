import re
from datetime import datetime
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation

from config import CHUNK_OVERLAP, CHUNK_SIZE, UPLOAD_DIR
from database import get_db


def extract_pdf_text(path: Path) -> tuple[str, int]:
    doc = fitz.open(path)
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n".join(pages), len(pages)


def extract_pptx_text(path: Path) -> tuple[str, int]:
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides.append("\n".join(parts))
    return "\n\n".join(slides), len(slides)


def guess_topic(filename: str, text: str) -> str:
    name = Path(filename).stem.replace("_", " ").replace("-", " ")
    if len(name) > 3:
        return name.title()[:80]
    first_line = text.strip().split("\n")[0][:80] if text.strip() else "General"
    return first_line or "General"


def chunk_text(text: str, title: str, topic: str, source: str) -> list[dict]:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []

    chunks = []
    start = 0
    idx = 0
    while start < len(text):
        end = start + CHUNK_SIZE
        piece = text[start:end].strip()
        if piece:
            chunks.append(
                {
                    "title": f"{title} — Part {idx + 1}",
                    "topic": topic,
                    "content": piece,
                    "source": source,
                }
            )
            idx += 1
        start = end - CHUNK_OVERLAP
        if start >= len(text):
            break

    return chunks


def process_upload(user_id: int, filename: str, file_bytes: bytes) -> dict:
    ext = Path(filename).suffix.lower()
    if ext not in (".pdf", ".pptx"):
        raise ValueError("Only PDF and PPTX files are supported")

    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    safe_name = re.sub(r"[^\w.\-]", "_", filename)
    save_path = UPLOAD_DIR / f"{user_id}_{safe_name}"
    save_path.write_bytes(file_bytes)

    if ext == ".pdf":
        text, page_count = extract_pdf_text(save_path)
        file_type = "pdf"
    else:
        text, page_count = extract_pptx_text(save_path)
        file_type = "pptx"

    title = Path(filename).stem
    topic = guess_topic(filename, text)
    source = filename
    chunks = chunk_text(text, title, topic, source)
    now = datetime.utcnow().isoformat()

    with get_db() as conn:
        cur = conn.execute(
            """INSERT INTO documents (user_id, filename, file_type, title, topic, page_count, chunk_count, created_at)
               VALUES (?, ?, ?, ?, ?, ?, ?, ?)""",
            (user_id, filename, file_type, title, topic, page_count, len(chunks), now),
        )
        doc_id = cur.lastrowid

        for ch in chunks:
            conn.execute(
                """INSERT INTO knowledge_chunks (document_id, user_id, title, topic, content, source, created_at)
                   VALUES (?, ?, ?, ?, ?, ?, ?)""",
                (doc_id, user_id, ch["title"], ch["topic"], ch["content"], ch["source"], now),
            )

    return {
        "document_id": doc_id,
        "title": title,
        "topic": topic,
        "page_count": page_count,
        "chunk_count": len(chunks),
    }
